"""Shared fixtures.

Every fixture builds **real files** — a genuine PDF via reportlab, a real DOCX,
PPTX and XLSX — rather than mocking the extraction libraries. Mocking pypdf would
only prove that the mock behaves like the mock; the failures worth catching are
in how these formats actually serialise.

Nothing here touches the network, and no test needs an API key.
"""

from __future__ import annotations

import io
from pathlib import Path

import pytest

from aidoctor.api.dependencies import build_container
from aidoctor.config.settings import Settings
from aidoctor.embeddings.base import build_embedder
from aidoctor.models.document import Document, Section, SourceType

HANDBOOK = {
    "Password reset": (
        "To reset your password open Settings and choose Reset Password. "
        "A confirmation email is sent to the registered address on file."
    ),
    "Billing": (
        "Invoices are issued monthly and include per-seat licence charges. "
        "Volume discounts apply above fifty seats."
    ),
    "Troubleshooting": (
        "ERR_LOCK_TIMEOUT indicates the work queue is saturated. "
        "Restart the worker pool to clear the condition."
    ),
}


@pytest.fixture()
def settings(tmp_path) -> Settings:
    """Isolated stack: temp SQLite file, embedded Qdrant, small vectors."""
    return Settings(
        database_url=f"sqlite:///{tmp_path / 'aidoctor.db'}",
        vector_backend="qdrant",
        embedding_dimensions=256,
        qdrant_collection="test",
    )


@pytest.fixture()
def container(settings):
    return build_container(settings)


@pytest.fixture()
def embedder():
    return build_embedder("hashing", 256)


@pytest.fixture()
def handbook() -> Document:
    return Document(
        doc_id="handbook",
        filename="handbook.md",
        source_type=SourceType.TEXT,
        sections=[
            Section(text=f"{label}\n{body}", label=label, ordinal=i)
            for i, (label, body) in enumerate(HANDBOOK.items())
        ],
    )


@pytest.fixture()
def pdf_file(tmp_path) -> Path:
    from reportlab.pdfgen import canvas

    path = tmp_path / "manual.pdf"
    pdf = canvas.Canvas(str(path))
    pdf.drawString(72, 720, "Licence keys are issued per seat and are non transferable.")
    pdf.showPage()
    pdf.drawString(72, 720, "Contact support for volume pricing above fifty seats.")
    pdf.save()
    return path


@pytest.fixture()
def docx_file(tmp_path) -> Path:
    from docx import Document as DocxDocument

    path = tmp_path / "handbook.docx"
    doc = DocxDocument()
    for label, body in HANDBOOK.items():
        doc.add_heading(label, 1)
        doc.add_paragraph(body)
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Plan"
    table.cell(0, 1).text = "Seats"
    table.cell(1, 0).text = "Enterprise"
    table.cell(1, 1).text = "500"
    doc.save(str(path))
    return path


@pytest.fixture()
def pptx_file(tmp_path) -> Path:
    from pptx import Presentation

    path = tmp_path / "deck.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Review"
    slide.placeholders[1].text = "Enterprise revenue grew across the licence base."
    slide.notes_slide.notes_text_frame.text = "Driven by three large renewals."
    deck.save(str(path))
    return path


@pytest.fixture()
def xlsx_file(tmp_path) -> Path:
    from openpyxl import Workbook

    path = tmp_path / "sales.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sales"
    sheet.append(["Region", "Amount", "Notes"])
    sheet.append(["EMEA", 1200, "renewal"])
    sheet.append([None, None, None])  # blank row: must be dropped
    sheet.append(["APAC", 900, ""])
    workbook.save(str(path))
    return path


@pytest.fixture()
def html_file(tmp_path) -> Path:
    path = tmp_path / "page.html"
    path.write_text(
        "<html><head><style>b{}</style><script>bad()</script></head><body>"
        "<h1>Setup</h1><p>Install the agent on each host.</p>"
        "<h2>Limits</h2><p>Ten seats maximum per licence.</p>"
        "</body></html>",
        encoding="utf-8",
    )
    return path


@pytest.fixture()
def text_file(tmp_path) -> Path:
    path = tmp_path / "notes.md"
    path.write_text("# Overview\nThe platform ingests documents.\n\n# Caveats\nNo OCR support.\n")
    return path


@pytest.fixture()
def all_files(pdf_file, docx_file, pptx_file, xlsx_file, html_file, text_file) -> list[Path]:
    return [pdf_file, docx_file, pptx_file, xlsx_file, html_file, text_file]


@pytest.fixture()
def client(settings):
    from fastapi.testclient import TestClient

    from aidoctor.main import create_app

    return TestClient(create_app(settings))


@pytest.fixture()
def docx_bytes(docx_file) -> bytes:
    return docx_file.read_bytes()


@pytest.fixture()
def upload(docx_bytes):
    def _upload(client, name: str = "handbook.docx", payload: bytes | None = None):
        return client.post(
            "/api/v1/ingest",
            files={
                "file": (name, payload if payload is not None else docx_bytes, "application/octet-stream")
            },
        )

    return _upload


def as_upload(path: Path) -> dict:
    return {"file": (path.name, io.BytesIO(path.read_bytes()), "application/octet-stream")}
