"""PPTX extraction, one section per slide.

Slide-level sections are the natural unit: a slide is already a self-contained
thought, and "slide 4" is a locator a reader can jump to.

Speaker notes are included when present. They frequently carry the substance
that the slide itself only gestures at, so dropping them loses the best content
in the deck.
"""

from __future__ import annotations

from pathlib import Path

from aidoctor.extractors.base import ExtractionError
from aidoctor.models.document import Document, Section, SourceType


class PptxExtractor:
    source_type = SourceType.PPTX
    extensions = (".pptx",)

    def extract(self, path: Path, doc_id: str) -> Document:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover - dependency guard
            raise ExtractionError("python-pptx is required to read .pptx files") from exc

        try:
            deck = Presentation(str(path))
        except Exception as exc:
            raise ExtractionError(f"Could not open {path.name}: {exc}") from exc

        sections: list[Section] = []
        for index, slide in enumerate(deck.slides, start=1):
            parts = [
                shape.text_frame.text.strip()
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    parts.append(f"Speaker notes: {notes}")
            if parts:
                sections.append(Section(text="\n".join(parts), label=f"slide {index}", ordinal=len(sections)))

        if not sections:
            raise ExtractionError(f"{path.name} contained no slide text")
        return Document(
            doc_id=doc_id,
            filename=path.name,
            source_type=self.source_type,
            sections=sections,
            metadata={"slides": str(len(deck.slides))},
        )
