#!/usr/bin/env python3
"""Regenerate ``examples/corpus`` from the definitions in this file.

Why a generator instead of hand-placed binaries: a .docx is a zip of XML, so a
reviewer cannot read a diff of one. Keeping the *text* here means the corpus is
reviewable, and the distractor design below is visible rather than folklore.

The corpus is deliberately adversarial. Retrieval over four short documents is
not a measurement — with fewer chunks than ``candidate_k`` every query returns
the whole index and Recall@k is 1.00 by arithmetic. So the documents here:

* **share vocabulary across files.** "password", "seat", "licence", "queue" and
  "restart" each appear in at least three documents, so no question is
  answerable by matching a single distinctive word.
* **include near-miss sections.** ``security-policy.docx`` describes password
  *rules*; only ``handbook.docx`` describes how to *reset* one. A retriever that
  keys on "password" alone ranks the wrong section first.
* **include sibling error codes.** ``ERR_LOCK_TIMEOUT``, ``ERR_LOCK_CONFLICT``
  and ``ERR_QUEUE_FULL`` all exist, so an exact-token match is not enough.
* **cover subjects with no answer**, to give the refusal metric something real
  to measure against.

Run: ``python scripts/make_corpus.py``
"""

from __future__ import annotations

from pathlib import Path

CORPUS = Path(__file__).resolve().parents[1] / "examples" / "corpus"

# --------------------------------------------------------------------------
# Markdown documents: "## " headings become section labels.
# --------------------------------------------------------------------------

RUNBOOK = """## Restart the worker pool
Drain the queue, then restart workers one at a time so in-flight jobs finish.
A full pool restart is the documented fix for a saturated work queue.

## Rollback
Re-deploy the previous image tag. Rollback does not revert database migrations.

## Scaling the pool
Add workers when queue depth stays above two hundred for five minutes.
Scaling out does not clear an existing lock; restart the pool for that.

## On-call handover
The outgoing engineer writes a handover note listing open incidents.

## Log locations
Worker logs are under /var/log/cobalt/worker. Retention is fourteen days.
"""

SUPPORT_FAQ = """## I cannot sign in
Check that the account is not suspended for non-payment before anything else.
If the password is simply forgotten, follow the reset steps in the handbook.

## Why did my invoice change
Invoices change when seat count changes mid-cycle. Charges are pro-rated.

## Can I move a licence to another person
A licence follows the seat, not the person. Reassign the seat instead.

## The app is slow after deployment
Slowness immediately after a deploy is usually queue backlog, not a lock.

## How do I add a teammate
An administrator assigns them a seat from the admin console.

## Do you offer a trial
A fourteen day trial is available and converts to a paid plan automatically.

## Where do I download invoices
Invoices are in the billing portal under Documents.

## Who do I contact about a security issue
Report security issues to the address in the security policy.
"""

POSTMORTEM = """## Summary
On 14 March the work queue saturated for forty minutes and jobs timed out.

## Impact
Roughly eight percent of scheduled jobs failed. No data was lost.

## Root cause
A slow downstream dependency held locks far longer than expected, so the queue
filled and workers reported timeouts.

## Resolution
The worker pool was restarted and the dependency timeout was lowered.

## Follow-up actions
Alert on queue depth, not only on error rate.
"""

GLOSSARY = """## Seat
A seat is a single assignable place on a plan. Billing counts seats, not people.

## Licence
A licence is the right to use the product, granted per seat.

## Work queue
The work queue holds background jobs waiting for a worker.

## Worker pool
The set of processes that consume the work queue.

## Lock
A lock prevents two workers from processing the same job.

## Pro-rating
Charging for part of a billing period when a seat changes mid-cycle.

## Tenant
An isolated customer account with its own documents and users.

## Retention
How long logs and backups are kept before deletion.
"""

RELEASE_NOTES = """## Version 4.2
Adds bulk seat assignment and a new billing portal export.

## Version 4.1
Reduces queue latency under load. Fixes a rare lock leak on worker shutdown.

## Version 4.0
Introduces the admin console. Removes the legacy licence key file.

## Deprecations
The legacy licence key file is no longer read as of 4.0.

## Upgrade notes
Upgrading from 3.x requires running the migration before starting workers.
"""

MARKDOWN = {
    "runbook.md": RUNBOOK,
    "support-faq.md": SUPPORT_FAQ,
    "postmortem.md": POSTMORTEM,
    "glossary.md": GLOSSARY,
    "release-notes.md": RELEASE_NOTES,
}

# --------------------------------------------------------------------------
# DOCX: headings become section labels.
# --------------------------------------------------------------------------

HANDBOOK = {
    "Password reset": (
        "To reset your password open Settings and choose Reset Password. "
        "A confirmation email is sent to the registered address on file."
    ),
    "Billing": (
        "Invoices are issued monthly and include per-seat licence charges. "
        "Volume discounts apply above fifty seats."
    ),
    "Troubleshooting": (
        "ERR_LOCK_TIMEOUT indicates the work queue is saturated. "
        "Restart the worker pool to clear the condition."
    ),
    "Leave policy": (
        "Annual leave is twenty five days plus public holidays. Requests go through the admin console."
    ),
    "Expenses": (
        "Submit expenses within thirty days with a receipt attached. Approval sits with the line manager."
    ),
    "Equipment": ("Laptops are replaced on a three year cycle. Report damage to the IT desk."),
    "Working hours": (
        "Core hours are ten to four in local time. Flexible start and finish times are supported."
    ),
    "Probation": ("The probation period is six months with a review at three."),
}

SECURITY_POLICY = {
    "Password requirements": (
        "Passwords must be at least twelve characters and are checked against a "
        "breach list. Reuse of the previous five passwords is blocked."
    ),
    "Password rotation": (
        "Scheduled password rotation is not required. Rotation is enforced only after a suspected compromise."
    ),
    "Multi-factor authentication": (
        "Multi-factor authentication is mandatory for administrators and "
        "strongly recommended for everyone else."
    ),
    "Access reviews": ("Seat and role assignments are reviewed every quarter."),
    "Reporting a vulnerability": (
        "Send vulnerability reports to security@example.invalid. Acknowledgement is within two working days."
    ),
    "Data retention": (
        "Customer documents are retained for the life of the tenant and deleted "
        "thirty days after termination."
    ),
}

DOCX_FILES = {"handbook.docx": HANDBOOK, "security-policy.docx": SECURITY_POLICY}

# --------------------------------------------------------------------------
# HTML: headings become section labels.
# --------------------------------------------------------------------------

API_ERRORS = [
    (
        "ERR_LOCK_TIMEOUT",
        "A worker waited too long for a lock. The work queue is saturated; restart the worker pool.",
    ),
    (
        "ERR_LOCK_CONFLICT",
        "Two workers claimed the same job. This is retried automatically and needs no action.",
    ),
    ("ERR_QUEUE_FULL", "The work queue rejected a new job because it is at capacity. Add workers."),
    ("ERR_SEAT_LIMIT", "The account has no free seat. Remove a seat assignment or buy more."),
    ("ERR_LICENCE_INVALID", "The licence could not be validated. Check the account is not suspended."),
    ("ERR_AUTH_FAILED", "Sign-in failed. The password may be wrong or the account suspended."),
    ("ERR_RATE_LIMITED", "Too many requests. Back off and retry with a delay."),
    ("ERR_DOC_TOO_LARGE", "The uploaded document exceeded the size limit."),
]

# --------------------------------------------------------------------------
# PPTX: one section per slide, labelled "slide N".
# --------------------------------------------------------------------------

ONBOARDING_SLIDES = [
    ("Welcome", "Your first week at a glance."),
    ("Getting an account", "An administrator assigns you a seat before day one."),
    ("Assigning seats", "Seats are assigned in the admin console and can be reassigned later."),
    ("Billing basics", "Seats are billed monthly and pro-rated when they change."),
    ("Getting help", "Support questions start in the FAQ, then the help desk."),
    ("Security", "Multi-factor authentication is set up on day one."),
]

# --------------------------------------------------------------------------
# PDF: one section per page.
# --------------------------------------------------------------------------

LICENCE_PAGES = [
    "Licence keys are issued per seat and are non transferable.",
    "The licence may be reassigned between users, but not shared concurrently.",
    "Audit rights allow verification of seat counts once per year.",
    "Termination for non-payment takes effect after thirty days notice.",
]

# --------------------------------------------------------------------------
# XLSX: one section per sheet.
# --------------------------------------------------------------------------

SALES_ROWS = [
    ("Region", "Amount", "Notes"),
    ("EMEA", 1200, "renewal"),
    ("APAC", 900, "new"),
]

PRICING_SHEETS = {
    "Tiers": [
        ("Tier", "Seats", "PricePerSeat"),
        ("Starter", "1-10", 12),
        ("Team", "11-50", 10),
        ("Business", "51+", 8),
    ],
    "Addons": [
        ("Addon", "Price", "Notes"),
        ("Extra storage", 40, "per terabyte per month"),
        ("Priority support", 250, "per account per month"),
    ],
}


def write_markdown() -> None:
    for name, body in MARKDOWN.items():
        (CORPUS / name).write_text(body, encoding="utf-8")


def write_docx() -> None:
    from docx import Document as Docx

    for name, sections in DOCX_FILES.items():
        doc = Docx()
        for heading, body in sections.items():
            doc.add_heading(heading, level=1)
            doc.add_paragraph(body)
        doc.save(CORPUS / name)


def write_html() -> None:
    parts = ["<html><body><h1>Error reference</h1>"]
    for code, body in API_ERRORS:
        parts.append(f"<h2>{code}</h2><p>{body}</p>")
    parts.append("</body></html>")
    (CORPUS / "api-errors.html").write_text("".join(parts), encoding="utf-8")


def write_pptx() -> None:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    for title, body in ONBOARDING_SLIDES:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    prs.save(CORPUS / "onboarding.pptx")


def write_pdf() -> None:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.pdfgen import canvas

    pdf = canvas.Canvas(str(CORPUS / "licence.pdf"), pagesize=LETTER)
    for text in LICENCE_PAGES:
        pdf.drawString(72, 720, text)
        pdf.showPage()
    pdf.save()


def write_xlsx() -> None:
    from openpyxl import Workbook

    book = Workbook()
    sheet = book.active
    sheet.title = "Sales"
    for row in SALES_ROWS:
        sheet.append(row)
    book.save(CORPUS / "sales.xlsx")

    book = Workbook()
    book.remove(book.active)
    for name, rows in PRICING_SHEETS.items():
        sheet = book.create_sheet(name)
        for row in rows:
            sheet.append(row)
    book.save(CORPUS / "pricing.xlsx")


def main() -> int:
    CORPUS.mkdir(parents=True, exist_ok=True)
    write_markdown()
    write_docx()
    write_html()
    write_pptx()
    write_pdf()
    write_xlsx()
    print(f"Wrote corpus to {CORPUS}")
    for path in sorted(CORPUS.iterdir()):
        if path.is_file():
            print(f"  {path.name:24} {path.stat().st_size:>7} bytes")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
